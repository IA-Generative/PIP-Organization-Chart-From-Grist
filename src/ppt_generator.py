from __future__ import annotations

from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
import os
import re
import sys
import threading
import time
import unicodedata
from typing import Dict, List, Sequence, Tuple

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .model_builder import BuiltModel, EpicModel, TeamModel

MARIANNE_FONT = "Marianne"
DEFAULT_SCW_BASE_URL = "https://api.scaleway.ai/a9158aac-8404-46ea-8bf5-1ca048cd6ab4/v1"
DEFAULT_SCW_MODEL = "mistral-small-3.2-24b-instruct-2506"
AI_CORRECTION_TAG = "(corrigé IA)"
AI_ROBOT_TAG = "(🤖)"
LONG_TEXT_THRESHOLD = 110
LONG_TEXT_MAX_WORDS = 52
LONG_TEXT_MAX_WORDS_FALLBACK = 104
TITLE_LONG_THRESHOLD = 56
_AI_REWRITE_CACHE: Dict[str, tuple[str, bool]] = {}
_LLM_DISABLED_REASON: str | None = None
_LLM_STATS = {"planned": 0, "calls": 0, "ok": 0, "empty": 0, "fallback": 0, "errors": 0}
_LLM_PROGRESS_TICK = 0
_LLM_STATS_LOCK = threading.Lock()
_AI_CACHE_LOCK = threading.Lock()
NO_FEATURES_LINE_PREFIX = "[NO_FEATURES]"
EPIC_HEADER_LINE_PREFIX = "[EPIC_HEADER]"
EPIC_HEADER_AMBITION_PREFIX = "[EPIC_HEADER_AMBITION]"
EPIC_HEADER_FEATURES_PREFIX = "[EPIC_HEADER_FEATURES]"
DARK_RED_HEX = "8B0000"
DARK_BLUE_HEX = "0070C0"
FINALITES_PAGE_MAX_LINES = 10


class _TerminalSpinner:
    def __init__(self, message: str, *, enabled: bool = True) -> None:
        self.message = message
        self._stop = threading.Event()
        self._thread: threading.Thread | None = None
        self._is_tty = sys.stdout.isatty() and enabled
        self._frames = ["|", "/", "-", "\\"]
        self._enabled = enabled

    def start(self) -> None:
        if not self._enabled:
            return
        if not self._is_tty:
            print(f"[LLM] {self.message}...")
            return

        def _run() -> None:
            idx = 0
            while not self._stop.is_set():
                frame = self._frames[idx % len(self._frames)]
                sys.stdout.write(f"\r[LLM] {frame} {self.message}...")
                sys.stdout.flush()
                time.sleep(0.12)
                idx += 1

        self._thread = threading.Thread(target=_run, daemon=True)
        self._thread.start()

    def stop(self, result: str) -> None:
        if not self._enabled:
            return
        self._stop.set()
        if self._thread:
            self._thread.join(timeout=0.2)
        if self._is_tty:
            # Clear current spinner line before writing final status.
            sys.stdout.write("\r\033[K")
            sys.stdout.write(f"[LLM] {result}\n")
            sys.stdout.flush()
        else:
            print(f"[LLM] {result}")


def _llm_log_mode() -> str:
    mode = os.getenv("LLM_LOG_MODE", "compact").strip().lower()
    if mode not in {"quiet", "compact", "verbose"}:
        return "compact"
    return mode


def _llm_log_event(status: str, context: str = "") -> None:
    global _LLM_PROGRESS_TICK
    mode = _llm_log_mode()
    if mode == "quiet":
        return

    with _LLM_STATS_LOCK:
        if status == "call":
            _LLM_STATS["calls"] += 1
        elif status == "ok":
            _LLM_STATS["ok"] += 1
        elif status == "empty":
            _LLM_STATS["empty"] += 1
            _LLM_STATS["fallback"] += 1
        elif status == "error":
            _LLM_STATS["errors"] += 1
            _LLM_STATS["fallback"] += 1

    if mode == "compact":
        with _LLM_STATS_LOCK:
            calls = _LLM_STATS["calls"]
            planned = _LLM_STATS["planned"]
            ok = _LLM_STATS["ok"]
            fb = _LLM_STATS["fallback"]
            done = ok + fb
        width = 16
        denom = planned if planned > 0 else calls
        ratio = (done / denom) if denom else 0.0
        filled = min(width, max(0, int(round(ratio * width))))
        bar = "#" * filled + "-" * (width - filled)
        spinner = ["|", "/", "-", "\\"][_LLM_PROGRESS_TICK % 4]
        _LLM_PROGRESS_TICK += 1
        msg = f"[LLM] {spinner} appels={calls} ok={ok} fallback={fb}"
        if context:
            msg += f" | bloc={context}"
        msg += f" | [{bar}] {done}/{denom if denom else 0}"
        if sys.stdout.isatty():
            sys.stdout.write("\r\033[K" + msg)
            sys.stdout.flush()
        else:
            if status in {"ok", "empty", "error"}:
                print(msg)
        return

    # verbose
    if context:
        print(f"[LLM] {status} | {context}")
    else:
        print(f"[LLM] {status}")


def _llm_log_summary() -> None:
    if not _llm_enabled():
        return
    mode = _llm_log_mode()
    if mode == "quiet":
        return
    if mode == "compact" and sys.stdout.isatty():
        sys.stdout.write("\r\033[K")
        sys.stdout.flush()
    with _LLM_STATS_LOCK:
        planned = _LLM_STATS["planned"]
        calls = _LLM_STATS["calls"]
        ok = _LLM_STATS["ok"]
        empty = _LLM_STATS["empty"]
        errors = _LLM_STATS["errors"]
        fallback = _LLM_STATS["fallback"]
    print(
        "[LLM] resume: "
        f"planifies={planned} appels={calls} ok={ok} "
        f"vide={empty} erreurs={errors} "
        f"fallback={fallback}"
    )


def _is_llm_candidate_line(line: str) -> bool:
    if line.startswith(EPIC_HEADER_LINE_PREFIX) or line.startswith(EPIC_HEADER_AMBITION_PREFIX) or line.startswith(EPIC_HEADER_FEATURES_PREFIX):
        return False
    raw = _normalize_line_text(line)
    return bool(raw) and len(raw) >= LONG_TEXT_THRESHOLD


def _estimate_planned_llm_calls(model: BuiltModel, frag_df) -> int:
    if not _llm_enabled():
        return 0

    total = 0

    def add_lines(lines: Sequence[str]) -> None:
        nonlocal total
        for line in lines:
            if _is_llm_candidate_line(line):
                total += 1

    # Slide 2/3/4 content lines (same logic as _fill_overview_slides)
    add_lines([
        f"PI: {model.pi}",
        "Objet: preparation PI planning",
        "Portee: equipes, epics, features",
    ])
    add_lines([
        f"Equipes: {model.stats.get('teams', 0)}",
        f"Epics: {model.stats.get('epics_total', 0)} (separees: {model.stats.get('epics_separate', 0)})",
        f"Features PI: {model.stats.get('features_pi', 0)}",
    ])
    add_lines([
        f"Agents references: {model.stats.get('personnes', 0)}",
        f"Agents avec affectation: {int(len(frag_df)) if frag_df is not None and not frag_df.empty else 0}",
        f"Nombre total d'affectations: {model.stats.get('affectations', 0)}",
    ])
    add_lines(["Nb_Epics >= 3 ou Nb_Equipes >= 2 ou charge totale > 100%"])
    if frag_df is not None and not frag_df.empty:
        low_df = frag_df.loc[frag_df["Total_Charge"] < 10.0].sort_values(
            ["Total_Charge", "Agent"], ascending=[True, True]
        )
        add_lines([f"{row['Agent']} - charge totale: {float(row['Total_Charge']):.1f}%" for _, row in low_df.head(30).iterrows()])
    else:
        add_lines(["Aucun agent en dessous de 10% de charge totale."])

    # Team groups
    for team in model.teams:
        add_lines([
            f"PM: {', '.join(team.pm_list) if team.pm_list else '-'}",
            f"PO: {', '.join(team.po_list) if team.po_list else '-'}",
            f"Nombre d'agents: {len(team.people_team)}",
            f"Epics associees: {len(team.epics)}",
        ])
        epic_without_features = sum(1 for e in team.epics if not e.features)
        add_lines([
            f"Epics sans feature PI: {epic_without_features}",
            "Dependances majeures: a completer selon contexte",
        ])
        finalites_lines = _finalites_for_team(team)
        ambitions_lines = _ambitions_for_team(team)
        features_lines = _features_for_team(team)[:18]
        add_lines(finalites_lines)
        add_lines(ambitions_lines)
        add_lines(features_lines)

    return total


def _rewrite_cache_key(text: str, max_words: int | None) -> str:
    return f"{max_words}|{_normalize_line_text(text)}"


def _collect_ppt_rewrite_jobs(model: BuiltModel, frag_df) -> List[Tuple[str, str, int | None]]:
    jobs: List[Tuple[str, str, int | None]] = []
    if not _llm_enabled():
        return jobs

    target_words = None

    def add_lines(title: str, lines: Sequence[str]) -> None:
        for line in lines:
            if _is_llm_candidate_line(line):
                jobs.append((title, line, target_words))

    # Overview slides
    add_lines("Informations cles PI", [
        f"PI: {model.pi}",
        "Objet: preparation PI planning",
        "Portee: equipes, epics, features",
    ])
    add_lines("Statistiques", [
        f"Equipes: {model.stats.get('teams', 0)}",
        f"Epics: {model.stats.get('epics_total', 0)} (separees: {model.stats.get('epics_separate', 0)})",
        f"Features PI: {model.stats.get('features_pi', 0)}",
    ])
    add_lines("Population d'agents", [
        f"Agents references: {model.stats.get('personnes', 0)}",
        f"Agents avec affectation: {int(len(frag_df)) if frag_df is not None and not frag_df.empty else 0}",
        f"Nombre total d'affectations: {model.stats.get('affectations', 0)}",
    ])
    add_lines("Criteres: Nb_Epics>=3 | Nb_Equipes>=2 | Charge_totale>100%", [])

    if frag_df is not None and not frag_df.empty:
        low_df = frag_df.loc[frag_df["Total_Charge"] < 10.0].sort_values(
            ["Total_Charge", "Agent"], ascending=[True, True]
        )
        add_lines(
            "Liste des agents concernes",
            [f"{row['Agent']} - charge totale: {float(row['Total_Charge']):.1f}%" for _, row in low_df.head(30).iterrows()],
        )
    else:
        add_lines("Liste des agents concernes", ["Aucun agent en dessous de 10% de charge totale."])

    # Team slides
    for team in model.teams:
        add_lines("Identite d'equipe", [
            f"PM: {', '.join(team.pm_list) if team.pm_list else '-'}",
            f"PO: {', '.join(team.po_list) if team.po_list else '-'}",
            f"Nombre d'agents: {len(team.people_team)}",
            f"Epics associees: {len(team.epics)}",
        ])
        epic_without_features = sum(1 for e in team.epics if not e.features)
        add_lines("Contexte", [
            f"Epics sans feature PI: {epic_without_features}",
            "Dependances majeures: a completer selon contexte",
        ])
        for page_idx, page in enumerate(_finalites_pages_for_team(team), start=1):
            add_lines("Finalites", page)
            if page_idx == 1:
                add_lines("Ambition du PIP", _ambitions_for_team(team))
        add_lines("Backlog features PI", _features_for_team(team)[:18])

    return jobs


def _prefetch_llm_rewrites(model: BuiltModel, frag_df) -> None:
    if not _llm_enabled():
        return

    jobs = _collect_ppt_rewrite_jobs(model, frag_df)
    if not jobs:
        return

    # De-duplicate identical rewrites to reduce API calls.
    unique_jobs: Dict[str, Tuple[str, str, int | None]] = {}
    for title, line, max_words in jobs:
        key = _rewrite_cache_key(line, max_words)
        if key not in unique_jobs:
            unique_jobs[key] = (title, line, max_words)

    max_workers = min(_ppt_llm_max_workers(), len(unique_jobs))
    if max_workers <= 1:
        for title, line, max_words in unique_jobs.values():
            _rewrite_text_with_fallback(line, max_words=max_words, context=title)
        return

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [
            executor.submit(_rewrite_text_with_fallback, line, max_words=max_words, context=title)
            for title, line, max_words in unique_jobs.values()
        ]
        for future in as_completed(futures):
            future.result()


def get_llm_status(deep_check: bool = False) -> tuple[bool, str]:
    if os.getenv("ENABLE_LLM", "0") != "1":
        return False, "flag --llm non active"
    api_key = os.getenv("SCW_SECRET_KEY_LLM")
    if not api_key:
        return False, "SCW_SECRET_KEY_LLM non definie"
    try:
        from openai import OpenAI
    except Exception as exc:
        return False, f"package openai indisponible ({type(exc).__name__})"
    if not deep_check:
        return True, "ok (precheck local)"
    base_url = os.getenv("SCW_LLM_BASE_URL", DEFAULT_SCW_BASE_URL)
    model = os.getenv("SCW_LLM_MODEL", DEFAULT_SCW_MODEL)
    try:
        client = OpenAI(base_url=base_url, api_key=api_key)
        client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": "ok"}],
            max_tokens=8,
            temperature=0,
            stream=False,
            response_format={"type": "text"},
        )
    except Exception as exc:
        msg = str(exc).replace("\n", " ").strip()
        if len(msg) > 140:
            msg = msg[:137] + "..."
        return False, f"echec appel API LLM ({type(exc).__name__}: {msg})"
    return True, "ok"


def _set_paragraph_font(paragraph, *, size_pt: int | None = None, bold: bool | None = None) -> None:
    if size_pt is not None:
        paragraph.font.size = Pt(size_pt)
    if bold is not None:
        paragraph.font.bold = bold
    paragraph.font.name = MARIANNE_FONT

    # Ensure font is applied to existing runs as well.
    for run in paragraph.runs:
        run.font.name = MARIANNE_FONT
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold


def _normalize_line_text(text: str) -> str:
    t = (text or "").strip()
    t = "".join(ch if (ch.isprintable() or ch in "\n\t") else " " for ch in t)
    t = re.sub(r"-{2,}", "-", t)
    t = re.sub(r"-\s+-", "-", t)
    t = re.sub(r"\s+", " ", t)
    t = re.sub(r"\s+([,.;:!?])", r"\1", t)
    return t


def _llm_enabled() -> bool:
    return os.getenv("ENABLE_LLM", "0") == "1"


def _local_rewrite_text(text: str) -> str:
    t = _normalize_line_text(text)
    if not t:
        return t
    # Keep bullets/labels as-is; only normalize prose-like lines.
    if t.startswith("- ") or ":" in t[:30]:
        return t
    if t and t[0].isalpha():
        t = t[0].upper() + t[1:]
    if t.endswith(":"):
        return t
    if t[-1] not in ".!?":
        t += "."
    return t


def _trim_to_words(text: str, max_words: int | None) -> str:
    if max_words is None:
        return _normalize_line_text(text)
    words = _normalize_line_text(text).split()
    if len(words) <= max_words:
        return " ".join(words)
    return " ".join(words[:max_words]).rstrip(",;:") + "..."


def _local_shorten_text(text: str, max_words: int | None) -> str:
    t = _normalize_line_text(text)
    if not t:
        return t
    if max_words is None:
        return _local_rewrite_text(t)
    parts = [p.strip() for p in re.split(r"(?<=[\.;!?])\s+", t) if p.strip()]
    if parts:
        first = parts[0]
        return _trim_to_words(first, max_words)
    return _trim_to_words(t, max_words)


def _llm_rewrite_text(text: str, *, max_words: int | None, context: str = "") -> str | None:
    global _LLM_DISABLED_REASON

    if os.getenv("ENABLE_LLM", "0") != "1":
        _LLM_DISABLED_REASON = "flag --llm non active"
        return None

    if _LLM_DISABLED_REASON:
        return None

    api_key = os.getenv("SCW_SECRET_KEY_LLM")
    if not api_key:
        _LLM_DISABLED_REASON = "SCW_SECRET_KEY_LLM non definie dans l'environnement"
        return None

    try:
        from openai import OpenAI
    except Exception as exc:
        _LLM_DISABLED_REASON = f"librairie openai indisponible ({type(exc).__name__})"
        return None

    base_url = os.getenv("SCW_LLM_BASE_URL", DEFAULT_SCW_BASE_URL)
    model = os.getenv("SCW_LLM_MODEL", DEFAULT_SCW_MODEL)
    debug_enabled = os.getenv("LLM_DEBUG", "").strip().lower() in {"1", "true", "yes", "on"}
    client = OpenAI(base_url=base_url, api_key=api_key)
    verbose_spinner = _llm_log_mode() == "verbose"
    spinner = _TerminalSpinner("appel LLM pour raccourcir un paragraphe", enabled=verbose_spinner)

    if max_words is None:
        prompt = (
            "Réécris ce texte en français correct (orthographe, grammaire, ponctuation), "
            "en conservant le sens métier sans troncature artificielle. "
            "Reste clair et lisible pour une slide. "
            "Ne pas inventer d'information. Une seule sortie texte, sans markdown, sans guillemets.\n\n"
            f"Texte:\n{text}"
        )
    else:
        prompt = (
            "Réécris ce texte en français correct (orthographe, grammaire, ponctuation), "
            "en conservant le sens métier et en le raccourcissant pour tenir sur une slide. "
            f"Maximum {max_words} mots. "
            "Ne pas inventer d'information. Une seule sortie texte, sans markdown, sans guillemets.\n\n"
            f"Texte:\n{text}"
        )

    def _extract_text(resp) -> str:
        try:
            if resp and getattr(resp, "choices", None):
                msg = resp.choices[0].message
                content = getattr(msg, "content", None)
                if isinstance(content, str):
                    return content.strip()
                if isinstance(content, list):
                    parts = []
                    for item in content:
                        if isinstance(item, str):
                            parts.append(item)
                        elif isinstance(item, dict):
                            txt = item.get("text")
                            if isinstance(txt, str):
                                parts.append(txt)
                            elif item.get("type") == "output_text" and isinstance(item.get("value"), str):
                                parts.append(item.get("value"))
                    return " ".join(parts).strip()
        except Exception:
            return ""
        return ""
    try:
        _llm_log_event("call", context=context)
        spinner.start()
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "Tu corriges et reformules avec précision, sans ajouter de contenu."},
                {"role": "user", "content": prompt},
            ],
            temperature=0.1,
            max_tokens=260,
            stream=False,
            response_format={"type": "text"},
        )
        content = _extract_text(resp)
        if debug_enabled:
            try:
                print(f"[LLM][DEBUG] model={model} base_url={base_url}")
                print(f"[LLM][DEBUG] finish_reason={resp.choices[0].finish_reason if resp and resp.choices else 'n/a'}")
                print(f"[LLM][DEBUG] content_len={len(content)}")
            except Exception:
                pass

        # Retry once without response_format if content is empty.
        if not content:
            if debug_enabled:
                print("[LLM][DEBUG] empty content on first attempt, retrying without response_format")
            resp2 = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": "Tu corriges et reformules avec précision, sans ajouter de contenu."},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.1,
                max_tokens=260,
                stream=False,
            )
            content = _extract_text(resp2)
            if debug_enabled:
                try:
                    print(f"[LLM][DEBUG] retry_finish_reason={resp2.choices[0].finish_reason if resp2 and resp2.choices else 'n/a'}")
                    print(f"[LLM][DEBUG] retry_content_len={len(content)}")
                except Exception:
                    pass

        content = _normalize_line_text(content)
        if not content:
            spinner.stop("reponse vide, fallback local")
            _llm_log_event("empty", context=context)
            return None
        spinner.stop("reponse recue")
        _llm_log_event("ok", context=context)
        return _trim_to_words(content, max_words)
    except Exception:
        _LLM_DISABLED_REASON = "echec appel API LLM (verifie token/base_url/model)"
        spinner.stop("echec appel LLM, fallback local")
        _llm_log_event("error", context=context)
        return None


def _rewrite_text_with_fallback(text: str, *, max_words: int | None = 26, context: str = "") -> tuple[str, bool]:
    raw = _normalize_line_text(text)
    if len(raw) < LONG_TEXT_THRESHOLD:
        return _local_rewrite_text(raw), False

    cache_key = f"{max_words}|{raw}"
    with _AI_CACHE_LOCK:
        cached = _AI_REWRITE_CACHE.get(cache_key)
    if cached:
        return cached

    llm = _llm_rewrite_text(raw, max_words=max_words, context=context)
    if llm:
        result = (llm, True)
        with _AI_CACHE_LOCK:
            _AI_REWRITE_CACHE[cache_key] = result
        return result

    result = (_local_shorten_text(raw, max_words=max_words), False)
    with _AI_CACHE_LOCK:
        _AI_REWRITE_CACHE[cache_key] = result
    return result


def _ppt_llm_max_workers() -> int:
    raw = os.getenv("LLM_PPT_MAX_WORKERS", "16").strip()
    try:
        value = int(raw)
    except Exception:
        value = 16
    return max(1, min(256, value))


def _split_long_text(text: str, max_len: int = 120) -> List[str]:
    t = _normalize_line_text(text)
    if len(t) <= max_len:
        return [t]

    # Prefer semantic splits first: sentence boundaries.
    parts = [p.strip() for p in re.split(r"(?<=[\.;!?])\s+", t) if p.strip()]
    if len(parts) > 1:
        return [f"- {p}" for p in parts]

    # Fallback: wrap a long sentence as one bullet + continuation lines.
    core = t[2:].strip() if t.startswith("- ") else t

    words = core.split(" ")
    chunks: List[str] = []
    current = ""
    for w in words:
        trial = f"{current} {w}".strip()
        if current and len(trial) > max_len:
            chunks.append(current)
            current = w
        else:
            current = trial
    if current:
        chunks.append(current)

    if not chunks:
        return [t]

    formatted: List[str] = []
    for idx, chunk in enumerate(chunks):
        if idx == 0:
            formatted.append(f"- {chunk}")
        else:
            formatted.append(f"  {chunk}")
    return formatted


def _write_rich_paragraph(paragraph, text: str, *, size_pt: int, default_bold: bool = False) -> None:
    paragraph.level = 0
    paragraph.line_spacing = 1.0
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.text = ""

    raw_line = text or ""
    color_hex = None
    force_bold = default_bold
    if raw_line.startswith(NO_FEATURES_LINE_PREFIX):
        raw_line = raw_line[len(NO_FEATURES_LINE_PREFIX) :].lstrip()
        color_hex = DARK_RED_HEX
    if raw_line.startswith(EPIC_HEADER_LINE_PREFIX):
        raw_line = raw_line[len(EPIC_HEADER_LINE_PREFIX) :].lstrip()
        color_hex = DARK_BLUE_HEX
        force_bold = True
    if raw_line.startswith(EPIC_HEADER_AMBITION_PREFIX):
        raw_line = raw_line[len(EPIC_HEADER_AMBITION_PREFIX) :].lstrip()
        color_hex = DARK_BLUE_HEX
        force_bold = True
        paragraph.space_before = Pt(6)
    if raw_line.startswith(EPIC_HEADER_FEATURES_PREFIX):
        raw_line = raw_line[len(EPIC_HEADER_FEATURES_PREFIX) :].lstrip()
        color_hex = DARK_BLUE_HEX
        force_bold = True
        paragraph.space_before = Pt(6)

    line = _normalize_line_text(raw_line)
    if not line:
        run = paragraph.add_run()
        run.text = "-"
        run.font.name = MARIANNE_FONT
        run.font.size = Pt(size_pt)
        run.font.bold = force_bold
        if color_hex:
            run.font.color.rgb = RGBColor.from_string(color_hex)
        return

    # If a short label is present (e.g. "KPIs: ..."), emphasize it.
    colon_idx = line.find(":")
    has_label = 1 < colon_idx <= 42

    if has_label and not force_bold:
        label = line[: colon_idx + 1]
        value = line[colon_idx + 1 :].strip()

        run_label = paragraph.add_run()
        run_label.text = label
        run_label.font.name = MARIANNE_FONT
        run_label.font.size = Pt(size_pt)
        run_label.font.bold = True
        if color_hex:
            run_label.font.color.rgb = RGBColor.from_string(color_hex)

        if value:
            run_value = paragraph.add_run()
            run_value.text = f" {value}"
            run_value.font.name = MARIANNE_FONT
            run_value.font.size = Pt(size_pt)
            run_value.font.bold = False
            if color_hex:
                run_value.font.color.rgb = RGBColor.from_string(color_hex)
        return

    run = paragraph.add_run()
    run.text = line
    run.font.name = MARIANNE_FONT
    run.font.size = Pt(size_pt)
    run.font.bold = force_bold
    if color_hex:
        run.font.color.rgb = RGBColor.from_string(color_hex)


def _write_block(shape, title: str, bullets: Sequence[str]) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.clear()

    p0 = tf.paragraphs[0]
    _write_rich_paragraph(p0, title, size_pt=18, default_bold=True)

    rewritten_rows: List[tuple[str, bool]] = [("", False)] * len(bullets)
    target_words = None if _llm_enabled() else LONG_TEXT_MAX_WORDS_FALLBACK

    def _is_passthrough(line: str) -> bool:
        return (
            line.startswith(EPIC_HEADER_LINE_PREFIX)
            or line.startswith(EPIC_HEADER_AMBITION_PREFIX)
            or line.startswith(EPIC_HEADER_FEATURES_PREFIX)
        )

    candidates: List[tuple[int, str]] = []
    for idx, line in enumerate(bullets):
        if _is_passthrough(line):
            rewritten_rows[idx] = (line, False)
        else:
            candidates.append((idx, line))

    if candidates and _llm_enabled():
        max_workers = min(_ppt_llm_max_workers(), len(candidates))
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(
                    _rewrite_text_with_fallback,
                    line,
                    max_words=target_words,
                    context=title,
                ): idx
                for idx, line in candidates
            }
            for future in as_completed(futures):
                idx = futures[future]
                rewritten_rows[idx] = future.result()
    else:
        for idx, line in candidates:
            rewritten_rows[idx] = _rewrite_text_with_fallback(
                line, max_words=target_words, context=title
            )

    for rewritten, used_ai in rewritten_rows:
        formatted_lines = _split_long_text(rewritten)
        if used_ai and formatted_lines:
            formatted_lines[-1] = f"{formatted_lines[-1]} {AI_ROBOT_TAG}"
        for formatted_line in formatted_lines:
            p = tf.add_paragraph()
            _write_rich_paragraph(p, formatted_line, size_pt=13)


def _write_title(slide, text: str) -> None:
    title_shape = slide.shapes.title
    if not title_shape:
        return
    title_shape.text = (text or "").upper()
    if title_shape.has_text_frame and title_shape.text_frame.paragraphs:
        # Titles must stay on horizontal lines; avoid aggressive auto-fit wrapping.
        title_shape.text_frame.word_wrap = False
        title_shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        title_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        title_shape.text_frame.margin_left = Inches(0.03)
        title_shape.text_frame.margin_right = Inches(0.03)
        title_shape.text_frame.margin_top = Inches(0.02)
        title_shape.text_frame.margin_bottom = Inches(0.02)
        title_len = len((text or "").strip())
        if title_len > 78:
            title_size_pt = 14
        elif title_len > TITLE_LONG_THRESHOLD:
            title_size_pt = 16
        else:
            title_size_pt = 20
        for p in title_shape.text_frame.paragraphs:
            _set_paragraph_font(p, size_pt=title_size_pt, bold=True)


def _content_text_shapes(slide) -> List:
    title = slide.shapes.title
    result = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if title is not None and shape == title:
            continue
        result.append(shape)
    return result


def _get_or_add_content_shape(slide, index: int, x: float, y: float, w: float, h: float):
    text_shapes = _content_text_shapes(slide)

    # Prefer the template text box whose geometry is closest to the requested one.
    if text_shapes:
        target_left = int(Inches(x))
        target_top = int(Inches(y))
        target_width = int(Inches(w))
        target_height = int(Inches(h))

        def _shape_score(shape) -> int:
            return (
                abs(int(shape.left) - target_left)
                + abs(int(shape.top) - target_top)
                + abs(int(shape.width) - target_width)
                + abs(int(shape.height) - target_height)
            )

        best_shape = min(text_shapes, key=_shape_score)
        # Keep backward compatibility for slides without a tuned template.
        if index < len(text_shapes):
            indexed_shape = text_shapes[index]
            if _shape_score(best_shape) <= _shape_score(indexed_shape):
                return best_shape
            return indexed_shape
        return best_shape

    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def _find_labeled_content_shape(slide, keyword: str):
    k = _normalize_text(keyword)
    if not k:
        return None
    for shape in _content_text_shapes(slide):
        try:
            text = _normalize_text(shape.text_frame.text or "")
        except Exception:
            text = ""
        if k in text:
            return shape
    return None


def _normalize_text(text: str) -> str:
    s = (text or "").strip().lower()
    s = unicodedata.normalize("NFKD", s)
    return "".join(ch for ch in s if not unicodedata.combining(ch))


def _pick_content_shape(slide, *, keyword: str, index: int, x: float, y: float, w: float, h: float):
    by_label = _find_labeled_content_shape(slide, keyword)
    if by_label is not None:
        return by_label
    return _get_or_add_content_shape(slide, index, x, y, w, h)


def _layout_by_name(prs: Presentation, name: str):
    for i in range(len(prs.slide_layouts)):
        if prs.slide_layouts[i].name == name:
            return prs.slide_layouts[i]
    return prs.slide_layouts[0]


def _epic_ambition(epic: EpicModel) -> str:
    parts = []
    if epic.description:
        parts.append(epic.description.strip())
    if epic.intention_pi:
        parts.append(f"Intention PI : {epic.intention_pi.strip()}")
    if epic.intention_next:
        parts.append(f"Prochain increment : {epic.intention_next.strip()}")
    text = " ".join(parts)
    if len(text) > 320 and not _llm_enabled():
        text = text[:317] + "..."
    return text or "-"


def _team_summary_lines(team: TeamModel) -> List[str]:
    lines: List[str] = [f"Equipe: {team.name}"]
    if team.mission_summary:
        lines.append(f"Mission: {team.mission_summary.strip()}")
    if team.next_increment_summary:
        lines.append(f"Prochain increment: {team.next_increment_summary.strip()}")
    if team.kpi_summary:
        lines.append(f"KPIs: {team.kpi_summary.strip()}")
    if team.summary_warning:
        lines.append(f"Alerte: {team.summary_warning.strip()}")
    return lines


def _features_for_team(team: TeamModel) -> List[str]:
    rows: List[str] = []
    for epic in team.epics:
        if not epic.features:
            rows.append(f"{EPIC_HEADER_FEATURES_PREFIX}{epic.name}")
            rows.append(f"{NO_FEATURES_LINE_PREFIX}  - aucune feature PI")
            continue
        rows.append(f"{EPIC_HEADER_FEATURES_PREFIX}{epic.name}")
        for feat in epic.features:
            if not _normalize_line_text(feat):
                rows.append(f"{NO_FEATURES_LINE_PREFIX}  - aucune feature PI")
                continue
            rows.append(f"  - {feat}")
    return rows or [f"{NO_FEATURES_LINE_PREFIX}Aucune feature PI identifiee pour cette equipe."]


def _set_table_cell_text(cell, text: str, *, bold: bool = False, size_pt: int = 11) -> None:
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.level = 0
    _set_paragraph_font(p, size_pt=size_pt, bold=bold)


def _add_fragmentation_table(slide, left, top, width, height, frag_df) -> None:
    headers = ["Agent", "Equipes", "Epics", "Affect.", "Charge %", "Score"]

    rows_data: List[List[str]] = []
    if frag_df is not None and not frag_df.empty:
        frag_focus = frag_df.loc[
            (frag_df["Nb_Epics"] >= 3) | (frag_df["Nb_Equipes"] >= 2) | (frag_df["Total_Charge"] > 100)
        ].copy()
        frag_focus = frag_focus.sort_values(
            ["Score_Fragmentation", "Total_Charge", "Nb_Epics", "Nb_Equipes"],
            ascending=[False, False, False, False],
        )
        for _, row in frag_focus.head(18).iterrows():
            rows_data.append(
                [
                    str(row["Agent"]),
                    str(int(row["Nb_Equipes"])),
                    str(int(row["Nb_Epics"])),
                    str(int(row["Nb_Affectations"])),
                    f"{float(row['Total_Charge']):.1f}",
                    str(int(row["Score_Fragmentation"])),
                ]
            )

    if not rows_data:
        rows_data = [["Aucun agent concerne", "-", "-", "-", "-", "-"]]

    rows = 1 + len(rows_data)
    cols = len(headers)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Ratios sum to 1.0 so the table uses the full allotted width (visually centered on slide).
    col_widths = [0.40, 0.10, 0.10, 0.10, 0.15, 0.15]
    for idx, ratio in enumerate(col_widths):
        table.columns[idx].width = int(width * ratio)

    for c, head in enumerate(headers):
        _set_table_cell_text(table.cell(0, c), head, bold=True, size_pt=11)

    for r, row_data in enumerate(rows_data, start=1):
        for c, value in enumerate(row_data):
            _set_table_cell_text(table.cell(r, c), value, size_pt=10)


def _build_team_member_rows(team: TeamModel) -> List[List[str]]:
    people = {p for p in team.people_team if _normalize_line_text(p) and _normalize_line_text(p) != "UNKNOWN"}
    people.update([p for p in team.pm_list if _normalize_line_text(p)])
    people.update([p for p in team.po_list if _normalize_line_text(p)])

    # Prefer direct team-level aggregation from Affectations table.
    total_charge_by_person: Dict[str, float] = dict(team.member_total_charge or {})
    roles_by_person: Dict[str, set[str]] = {
        person: set(roles) for person, roles in (team.member_roles or {}).items()
    }
    people.update([p for p in total_charge_by_person.keys() if _normalize_line_text(p)])

    # Fallback/merge from epic assignments if needed.
    for epic in team.epics:
        for a in epic.assignments:
            person = _normalize_line_text(a.person)
            if not person or person == "UNKNOWN":
                continue
            if person not in people:
                people.add(person)
            if person not in total_charge_by_person:
                total_charge_by_person[person] = float(a.charge or 0.0)
            role = _normalize_line_text(a.role)
            if role:
                roles_by_person.setdefault(person, set()).add(role)

    rows: List[List[str]] = []
    for person in sorted(people):
        roles = set(roles_by_person.get(person, set()))
        if person in team.pm_list:
            roles.add("PM")
        if person in team.po_list:
            roles.add("PO")
        role_label = ", ".join(sorted(roles)) if roles else "Membre"
        charge = total_charge_by_person.get(person, 0.0)
        charge_display = round(charge, 1)
        if charge_display <= 0.0:
            continue
        rows.append([person, role_label, f"{charge_display:.1f}%"])

    rows.sort(key=lambda r: float(r[2].replace("%", "")), reverse=True)
    return rows


def _add_team_members_table(slide, left, top, width, height, team: TeamModel) -> None:
    headers = ["Membre", "Qualite", "Affectation %"]
    rows_data = _build_team_member_rows(team)
    if not rows_data:
        rows_data = [["Aucun membre detecte", "-", "0.0%"]]

    # Keep table readable in fixed block height.
    rows_data = rows_data[:18]
    rows = 1 + len(rows_data)
    cols = len(headers)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    col_widths = [0.50, 0.25, 0.25]
    for idx, ratio in enumerate(col_widths):
        table.columns[idx].width = int(width * ratio)

    for c, head in enumerate(headers):
        _set_table_cell_text(table.cell(0, c), head, bold=True, size_pt=11)
    for r, row_data in enumerate(rows_data, start=1):
        for c, value in enumerate(row_data):
            _set_table_cell_text(table.cell(r, c), value, size_pt=10)


def _add_empty_team_members_table(slide, left, top, width, height) -> None:
    table_shape = slide.shapes.add_table(2, 3, left, top, width, height)
    table = table_shape.table
    headers = ["Membre", "Qualite", "Affectation %"]
    col_widths = [0.50, 0.25, 0.25]
    for idx, ratio in enumerate(col_widths):
        table.columns[idx].width = int(width * ratio)
    for c, h in enumerate(headers):
        _set_table_cell_text(table.cell(0, c), h, bold=True, size_pt=11)
    _set_table_cell_text(table.cell(1, 0), "Aucun membre detecte", size_pt=10)
    _set_table_cell_text(table.cell(1, 1), "-", size_pt=10)
    _set_table_cell_text(table.cell(1, 2), "0.0%", size_pt=10)


def _finalites_for_team(team: TeamModel) -> List[str]:
    rows: List[str] = []
    for epic in team.epics:
        desc = epic.description.strip() if epic.description else "sans description"
        rows.append(f"{EPIC_HEADER_LINE_PREFIX}{epic.name}:")
        rows.extend(_split_long_text(desc, max_len=105))
    return rows or ["Aucune finalite renseignee."]


def _finalites_pages_for_team(team: TeamModel, max_lines: int = FINALITES_PAGE_MAX_LINES) -> List[List[str]]:
    sections: List[List[str]] = []
    for epic in team.epics:
        section: List[str] = [f"{EPIC_HEADER_LINE_PREFIX}{epic.name}:"]
        desc = epic.description.strip() if epic.description else "sans description"
        section.extend(_split_long_text(desc, max_len=105))
        sections.append(section)

    if not sections:
        return [["Aucune finalite renseignee."]]

    pages: List[List[str]] = []
    current: List[str] = []
    current_count = 0
    for section in sections:
        sec_count = len(section)
        if current and current_count + sec_count > max_lines:
            pages.append(current)
            current = []
            current_count = 0

        if sec_count > max_lines:
            # Very large section: split but keep header at the start of each chunk.
            header = section[0]
            body = section[1:]
            chunk: List[str] = [header]
            chunk_count = 1
            for line in body:
                if chunk_count + 1 > max_lines:
                    pages.append(chunk)
                    chunk = [header, line]
                    chunk_count = 2
                else:
                    chunk.append(line)
                    chunk_count += 1
            if chunk:
                current = chunk
                current_count = len(chunk)
            continue

        current.extend(section)
        current_count += sec_count

    if current:
        pages.append(current)
    return pages


def _ambitions_for_team(team: TeamModel) -> List[str]:
    rows: List[str] = []
    for epic in team.epics:
        ambition = _epic_ambition(epic)
        rows.append(f"{EPIC_HEADER_AMBITION_PREFIX}{epic.name}:")
        rows.extend(_split_long_text(ambition, max_len=105))
    return rows or ["Aucune ambition PI renseignee."]


def _ensure_base_slides(prs: Presentation) -> None:
    # Template should contain 7 slides; if not, add missing ones.
    while len(prs.slides) < 7:
        if len(prs.slides) == 0:
            prs.slides.add_slide(_layout_by_name(prs, "Title Slide"))
        else:
            prs.slides.add_slide(_layout_by_name(prs, "Title Only"))


def _fill_overview_slides(prs: Presentation, model: BuiltModel, frag_df) -> None:
    # Slide 1
    s1 = prs.slides[0]
    _write_title(s1, f"PI Planning - Synthese generale ({model.pi})")
    if len(s1.placeholders) > 1:
        subtitle = s1.placeholders[1]
        subtitle.text = "Synthese automatique des equipes, epics, features et affectations"
        if subtitle.has_text_frame:
            subtitle.text_frame.word_wrap = True
            subtitle.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            subtitle.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            subtitle.text_frame.margin_left = Inches(0.03)
            subtitle.text_frame.margin_right = Inches(0.03)
            subtitle.text_frame.margin_top = Inches(0.02)
            subtitle.text_frame.margin_bottom = Inches(0.02)
            for p in subtitle.text_frame.paragraphs:
                _set_paragraph_font(p, size_pt=14)

    # Slide 2
    s2 = prs.slides[1]
    _write_title(s2, "Vue d'ensemble du PI a venir")

    total_people = int(model.stats.get("personnes", 0))
    active_people = int(len(frag_df)) if frag_df is not None and not frag_df.empty else 0
    total_assign = int(model.stats.get("affectations", 0))

    b1 = _pick_content_shape(s2, keyword="information", index=0, x=0.8, y=1.4, w=4.6, h=2.2)
    _write_block(
        b1,
        "Informations cles PI",
        [
            f"PI: {model.pi}",
            "Objet: preparation PI planning",
            "Portee: equipes, epics, features",
        ],
    )

    b2 = _pick_content_shape(s2, keyword="statist", index=1, x=5.0, y=1.4, w=4.3, h=2.2)
    _write_block(
        b2,
        "Statistiques",
        [
            f"Equipes: {model.stats.get('teams', 0)}",
            f"Epics: {model.stats.get('epics_total', 0)} (separees: {model.stats.get('epics_separate', 0)})",
            f"Features PI: {model.stats.get('features_pi', 0)}",
        ],
    )

    b3 = _pick_content_shape(s2, keyword="population", index=2, x=0.8, y=4.0, w=12.2, h=2.4)
    _write_block(
        b3,
        "Population d'agents",
        [
            f"Agents references: {total_people}",
            f"Agents avec affectation: {active_people}",
            f"Nombre total d'affectations: {total_assign}",
        ],
    )

    # Slide 3
    s3 = prs.slides[2]
    _write_title(s3, "Agents avec fragmentation d'affectation")
    b31 = _pick_content_shape(s3, keyword="crit", index=0, x=0.8, y=1.4, w=8.5, h=1.8)
    _write_block(
        b31,
        "Criteres: Nb_Epics>=3 | Nb_Equipes>=2 | Charge_totale>100%",
        [],
    )

    b32 = _pick_content_shape(s3, keyword="liste des agents", index=1, x=0.8, y=3.3, w=12.2, h=3.4)
    # Keep the template block as section title, then render tabular data in the same area.
    _write_block(b32, "Liste des agents concernes", [])
    _add_fragmentation_table(
        s3,
        left=b32.left,
        top=b32.top + int(Inches(0.42)),
        width=b32.width,
        height=max(int(Inches(0.8)), b32.height - int(Inches(0.45))),
        frag_df=frag_df,
    )

    # Slide 4
    s4 = prs.slides[3]
    _write_title(s4, "Agents avec faible affectation (<10%)")
    low_rows: List[str] = []
    if frag_df is not None and not frag_df.empty:
        low_df = frag_df.loc[frag_df["Total_Charge"] < 10.0].sort_values(
            ["Total_Charge", "Agent"], ascending=[True, True]
        )
        for _, row in low_df.head(30).iterrows():
            low_rows.append(f"{row['Agent']} - charge totale: {float(row['Total_Charge']):.1f}%")
    if not low_rows:
        low_rows = ["Aucun agent en dessous de 10% de charge totale."]

    b41 = _pick_content_shape(s4, keyword="perim", index=0, x=0.8, y=1.4, w=8.5, h=1.8)
    _write_block(
        b41,
        "Perimetre",
        ["Agents dont la charge totale est strictement inferieure a 10%"],
    )

    b42 = _pick_content_shape(s4, keyword="liste des agents", index=1, x=0.8, y=3.3, w=12.2, h=3.4)
    _write_block(b42, "Liste des agents concernes", low_rows)


def _fill_team_slide(s_team, team: TeamModel) -> None:
    # Slide team
    _write_title(s_team, f"Equipe - {team.name}")
    epic_without_features = sum(1 for e in team.epics if not e.features)

    b1 = _pick_content_shape(s_team, keyword="identite", index=0, x=0.8, y=1.4, w=4.6, h=2.2)
    _write_block(
        b1,
        "Identite d'equipe",
        [
            f"PM: {', '.join(team.pm_list) if team.pm_list else '-'}",
            f"PO: {', '.join(team.po_list) if team.po_list else '-'}",
            f"Nombre d'agents: {len(team.people_team)}",
            f"Epics associees: {len(team.epics)}",
        ],
    )

    b2 = _pick_content_shape(s_team, keyword="contexte", index=1, x=5.0, y=1.4, w=4.3, h=2.2)
    _write_block(
        b2,
        "Contexte",
        [
            f"Epics sans feature PI: {epic_without_features}",
            "Dependances majeures: a completer selon contexte",
        ],
    )

    b3 = _pick_content_shape(s_team, keyword="synthese", index=2, x=0.8, y=4.0, w=12.2, h=2.4)
    _write_block(b3, "Synthese equipe", [])
    _add_team_members_table(
        s_team,
        left=b3.left,
        top=b3.top + int(Inches(0.42)),
        width=b3.width,
        height=max(int(Inches(0.8)), b3.height - int(Inches(0.45))),
        team=team,
    )

def _fill_finalites_ambition_slide(s_ambition, team: TeamModel, finalites_lines: List[str], page_index: int, page_total: int) -> None:
    suffix = f" ({page_index}/{page_total})" if page_total > 1 else ""
    _write_title(s_ambition, f"Finalites et ambition du PIP - {team.name}{suffix}")
    a1 = _pick_content_shape(s_ambition, keyword="finalit", index=0, x=0.8, y=1.4, w=12.2, h=2.4)
    if not _llm_enabled():
        finalites_lines = finalites_lines[:8]
    _write_block(a1, "Finalites", finalites_lines)

    a2 = _pick_content_shape(s_ambition, keyword="ambition", index=1, x=0.8, y=4.0, w=12.2, h=2.4)
    ambitions_lines: List[str] = []
    if page_index == 1:
        ambitions_lines = _ambitions_for_team(team)
        if not _llm_enabled():
            ambitions_lines = ambitions_lines[:8]
    _write_block(a2, "Ambition du PIP", ambitions_lines)

def _fill_features_slide(s_features, team: TeamModel) -> None:
    _write_title(s_features, f"Features - {team.name}")
    f1 = _pick_content_shape(s_features, keyword="backlog", index=0, x=0.8, y=1.4, w=12.2, h=4.8)
    _write_block(f1, "Backlog features PI", _features_for_team(team)[:18])


def _enable_wrap_for_all_text(prs: Presentation) -> None:
    for slide in prs.slides:
        title_shape = slide.shapes.title
        title_shape_id = title_shape.shape_id if title_shape is not None else None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if title_shape_id is not None and shape.shape_id == title_shape_id:
                # Keep title behavior configured in _write_title.
                continue
            shape.text_frame.word_wrap = True
            shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP


def generate_ppt(model: BuiltModel, frag_kpis: Dict[str, int], out_path: str, frag_df=None) -> None:
    del frag_kpis  # Detailed values come from model + frag_df in this template-driven generator.
    global _LLM_PROGRESS_TICK
    _LLM_STATS["planned"] = _estimate_planned_llm_calls(model, frag_df)
    _LLM_STATS["calls"] = 0
    _LLM_STATS["ok"] = 0
    _LLM_STATS["empty"] = 0
    _LLM_STATS["fallback"] = 0
    _LLM_STATS["errors"] = 0
    _LLM_PROGRESS_TICK = 0
    if _llm_enabled() and _llm_log_mode() in {"compact", "verbose"}:
        print(f"[LLM] appels planifies: {_LLM_STATS['planned']}")

    # Precompute rewrites globally to maximize parallelism across all PPT blocks.
    _prefetch_llm_rewrites(model, frag_df)

    template_path = Path(__file__).resolve().parents[1] / "data" / "template.ppt.pptx"
    prs = Presentation(str(template_path)) if template_path.exists() else Presentation()

    _ensure_base_slides(prs)
    _fill_overview_slides(prs, model, frag_df)

    title_only_layout = _layout_by_name(prs, "Title Only")

    if model.teams:
        # First team uses pre-existing template slide for team summary.
        first_team = model.teams[0]
        _fill_team_slide(prs.slides[4], first_team)

        first_pages = _finalites_pages_for_team(first_team)
        _fill_finalites_ambition_slide(
            prs.slides[5],
            first_team,
            first_pages[0],
            page_index=1,
            page_total=len(first_pages),
        )
        if len(first_pages) == 1:
            _fill_features_slide(prs.slides[6], first_team)
        else:
            _fill_finalites_ambition_slide(
                prs.slides[6],
                first_team,
                first_pages[1],
                page_index=2,
                page_total=len(first_pages),
            )
            for idx in range(2, len(first_pages)):
                s_more = prs.slides.add_slide(title_only_layout)
                _fill_finalites_ambition_slide(
                    s_more,
                    first_team,
                    first_pages[idx],
                    page_index=idx + 1,
                    page_total=len(first_pages),
                )
            s_features_first = prs.slides.add_slide(title_only_layout)
            _fill_features_slide(s_features_first, first_team)

        for team in model.teams[1:]:
            s_team = prs.slides.add_slide(title_only_layout)
            _fill_team_slide(s_team, team)

            pages = _finalites_pages_for_team(team)
            for idx, page_lines in enumerate(pages, start=1):
                s_ambition = prs.slides.add_slide(title_only_layout)
                _fill_finalites_ambition_slide(
                    s_ambition,
                    team,
                    page_lines,
                    page_index=idx,
                    page_total=len(pages),
                )

            s_features = prs.slides.add_slide(title_only_layout)
            _fill_features_slide(s_features, team)
    else:
        _write_title(prs.slides[4], "Equipe - Aucune equipe detectee")
        _write_title(prs.slides[5], "Finalites et ambition du PIP - Aucune donnee")
        _write_title(prs.slides[6], "Features - Aucune donnee")
        empty_team_block = _pick_content_shape(prs.slides[4], keyword="synthese", index=2, x=0.8, y=4.0, w=12.2, h=2.4)
        _write_block(empty_team_block, "Synthese equipe", [])
        _add_empty_team_members_table(
            prs.slides[4],
            left=empty_team_block.left,
            top=empty_team_block.top + int(Inches(0.42)),
            width=empty_team_block.width,
            height=max(int(Inches(0.8)), empty_team_block.height - int(Inches(0.45))),
        )

    _enable_wrap_for_all_text(prs)
    prs.save(out_path)
    _llm_log_summary()
